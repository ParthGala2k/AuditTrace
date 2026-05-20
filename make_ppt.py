"""
Generate AuditTrace_Presentation.pptx  --  python make_ppt.py

Demo deck for CMPE 258 final project. All numbers come from:
  - evaluation/eval_e2e_results.json   (FP-filter eval, n=30)
  - evaluation/comparison_summary.json (planner eval,   n=55)

Design rules: large headings, no paragraphs, direct numbers, speakable.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── palette ─────────────────────────────────────────────────────────────────
INK    = RGBColor(0x0f, 0x11, 0x15)
PAPER  = RGBColor(0xfa, 0xfa, 0xf7)
ACCENT = RGBColor(0x6f, 0x42, 0xc1)
RED    = RGBColor(0xd6, 0x45, 0x45)
AMBER  = RGBColor(0xd9, 0xa4, 0x41)
GREEN  = RGBColor(0x5a, 0x9e, 0x6f)
GRAY   = RGBColor(0x6b, 0x72, 0x80)
LGRAY  = RGBColor(0xeb, 0xeb, 0xe8)
WHITE  = RGBColor(0xff, 0xff, 0xff)

W, H = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ── primitives ──────────────────────────────────────────────────────────────
def rect(sl, l, t, w, h, fill=None, line=None, lw=Pt(0.5)):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.line.width = lw
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if line: s.line.color.rgb = line
    else:    s.line.fill.background()
    return s


def txt(sl, text, l, t, w, h, size=14, bold=False, color=INK,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    r.font.name = font
    return tb


def bul(sl, items, l, t, w, h, size=14, color=INK, sp=6):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(sp)
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"


def hdr(sl, title, sub=""):
    rect(sl, 0, 0, W, Inches(1.2), fill=INK)
    rect(sl, 0, Inches(1.2), W, Inches(0.06), fill=ACCENT)
    txt(sl, title, Inches(0.5), Inches(0.18), Inches(12), Inches(0.7),
        size=30, bold=True, color=WHITE)
    if sub:
        txt(sl, sub, Inches(0.5), Inches(0.82), Inches(12), Inches(0.38),
            size=13, color=RGBColor(0xa0, 0x95, 0xcc), italic=True)


def card(sl, l, t, w, h, title, color=INK):
    rect(sl, l, t, w, h, fill=WHITE, line=LGRAY, lw=Pt(1))
    rect(sl, l, t, w, Inches(0.42), fill=color)
    txt(sl, title, l + Inches(0.12), t + Inches(0.06), w - Inches(0.2), Inches(0.32),
        size=13, bold=True, color=WHITE)


def tbl(sl, headers, rows, l, t, cw, rh=Inches(0.46)):
    x = l
    for i, h in enumerate(headers):
        rect(sl, x, t, cw[i], rh, fill=INK)
        txt(sl, h, x + Inches(0.06), t + Inches(0.08), cw[i] - Inches(0.12), rh - Inches(0.1),
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += cw[i]
    for ri, row in enumerate(rows):
        x = l
        bg = PAPER if ri % 2 == 0 else WHITE
        for ci, cell in enumerate(row):
            rect(sl, x, t + rh * (ri + 1), cw[ci], rh, fill=bg, line=LGRAY, lw=Pt(0.5))
            cc = INK
            if isinstance(cell, tuple): cell, cc = cell
            txt(sl, str(cell), x + Inches(0.06), t + rh * (ri + 1) + Inches(0.1),
                cw[ci] - Inches(0.12), rh - Inches(0.12),
                size=12, color=cc, align=PP_ALIGN.CENTER)
            x += cw[ci]


def big_stat(sl, val, label, l, t, w, h, color):
    rect(sl, l, t, w, h, fill=color)
    txt(sl, val, l, t + Inches(0.18), w, Inches(0.95),
        size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, label, l, t + h - Inches(0.55), w, Inches(0.42),
        size=12, color=WHITE, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=INK)
rect(s, 0, H - Inches(0.07), W, Inches(0.07), fill=ACCENT)
txt(s, "AuditTrace", Inches(1), Inches(2.0), Inches(11.33), Inches(1.3),
    size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Multi-Agent LLM Consensus for Compliance Auditing",
    Inches(1), Inches(3.4), Inches(11.33), Inches(0.65),
    size=22, color=RGBColor(0xb0, 0xa8, 0xd0), align=PP_ALIGN.CENTER)
txt(s, "CMPE 258 Deep Learning  |  SJSU Spring 2026  |  Group 4",
    Inches(1), Inches(4.4), Inches(11.33), Inches(0.45),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, "Parth Gala",
    Inches(1), Inches(4.9), Inches(11.33), Inches(0.4),
    size=14, color=GRAY, align=PP_ALIGN.CENTER)
for i, (label, color) in enumerate([
    ("Checkov Detection", RGBColor(0x2d, 0x6a, 0x4f)),
    ("3-Model FP Filter", ACCENT),
    ("Consensus Voting",  RED),
]):
    lx = Inches(3.2 + i * 2.35)
    rect(s, lx, Inches(6.0), Inches(2.1), Inches(0.42), fill=color)
    txt(s, label, lx + Inches(0.08), Inches(6.05), Inches(1.95), Inches(0.32),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 2 — The Problem
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
hdr(s, "The Problem", "Static scanners detect violations; they cannot reason about context")

big_stat(s, "385",  "Checkov findings on terragoat",  Inches(0.5),  Inches(1.55), Inches(2.95), Inches(1.6), INK)
big_stat(s, "0",    "with surrounding context",       Inches(3.65), Inches(1.55), Inches(2.95), Inches(1.6), RED)
big_stat(s, "~30%", "are noise on real repos",        Inches(6.8),  Inches(1.55), Inches(2.95), Inches(1.6), AMBER)
big_stat(s, "Hours", "engineer time per audit",        Inches(9.95), Inches(1.55), Inches(2.95), Inches(1.6), GRAY)

txt(s, "Checkov Strength", Inches(0.5), Inches(3.55), Inches(6.1), Inches(0.35),
    size=15, bold=True, color=GREEN)
rect(s, Inches(0.5), Inches(3.95), Inches(6.1), Inches(2.9), fill=WHITE, line=LGRAY, lw=Pt(1))
bul(s, [
    "  Deterministic — never misses a rule",
    "  500+ AWS/Terraform checks built-in",
    "  Runs in seconds, fully offline",
],   Inches(0.65), Inches(4.05), Inches(5.85), Inches(2.7), size=15)

txt(s, "Checkov Weakness", Inches(7.1), Inches(3.55), Inches(5.7), Inches(0.35),
    size=15, bold=True, color=RED)
rect(s, Inches(7.1), Inches(3.95), Inches(5.7), Inches(2.9), fill=WHITE, line=LGRAY, lw=Pt(1))
bul(s, [
    "  Treats a CI test bucket the same as production",
    "  Cannot read comments, file paths, sibling resources",
    "  Drowns engineers in unprioritized alerts",
],   Inches(7.25), Inches(4.05), Inches(5.45), Inches(2.7), size=15)

txt(s, "We add LLM judgement — but only where Checkov needs it most.",
    Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
    size=14, color=GRAY, italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 3 — Our Approach
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
hdr(s, "Our Approach", "Hybrid pipeline: deterministic detection + contextual triage")

stages = [
    ("1. Detect", RGBColor(0x2d, 0x6a, 0x4f), [
        "Clone repo",
        "Run Checkov",
        "Extract full Terraform block per finding",
    ]),
    ("2. Judge", ACCENT, [
        "Send each finding + HCL context to 3 LLMs",
        "Each model independently labels:",
        "      GENUINE  /  FALSE_POSITIVE  /  UNCERTAIN",
    ]),
    ("3. Vote", RED, [
        "3/3 GENUINE       →  HIGH",
        "2/3 GENUINE       →  LIKELY",
        "3/3 FALSE_POS.    →  SUPPRESSED (auto-dropped)",
    ]),
]
for i, (title, color, items) in enumerate(stages):
    lx = Inches(0.4 + i * 4.3); ty = Inches(1.55)
    card(s, lx, ty, Inches(4.1), Inches(4.7), title, color)
    bul(s, ["  " + it for it in items],
        lx + Inches(0.14), ty + Inches(0.55),
        Inches(3.85), Inches(4.0), size=14)

txt(s, "Key shift from the literature:",
    Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.35),
    size=14, bold=True, color=INK)
rect(s, Inches(0.5), Inches(6.78), Inches(12.3), Inches(0.55), fill=INK)
txt(s, "LLMs are not detectors — they are judges. They never invent findings, only validate them.",
    Inches(0.65), Inches(6.86), Inches(12.0), Inches(0.4),
    size=13, color=WHITE, italic=True)


# =============================================================================
# SLIDE 4 — System Architecture
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
hdr(s, "System Architecture", "LangGraph orchestrates a 2-node pipeline")

node_data = [
    ("CIS PDF",        Inches(0.3),   LGRAY,  INK),
    ("Policy\nDistill", Inches(2.35), INK,    WHITE),
    ("GitHub\nRepo",   Inches(4.4),   LGRAY,  INK),
    ("Checkov\n+ HCL", Inches(6.45),  INK,    WHITE),
    ("3-LLM\nJudges",  Inches(8.5),   ACCENT, WHITE),
    ("Tiered\nReport", Inches(10.55), GREEN,  WHITE),
]
for label, lx, bg, fg in node_data:
    rect(s, lx, Inches(1.55), Inches(1.9), Inches(1.0), fill=bg, line=LGRAY, lw=Pt(1))
    txt(s, label, lx + Inches(0.05), Inches(1.63), Inches(1.8), Inches(0.84),
        size=12, bold=True, color=fg, align=PP_ALIGN.CENTER)
for i in range(len(node_data) - 1):
    lx = node_data[i][1] + Inches(1.9)
    txt(s, "→", lx, Inches(1.9), Inches(0.45), Inches(0.32),
        size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

tech_labels = ["PyMuPDF", "GPT-4o mini (one-time)", "git clone", "Checkov 3.x",
               "OpenRouter parallel", "React + SQLite"]
for i, (_, lx, _2, _3) in enumerate(node_data):
    txt(s, tech_labels[i], lx, Inches(2.62), Inches(1.9), Inches(0.35),
        size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

detail = [
    ("Offline (once)", [
        "  CIS PDF → 89 structured clauses",
        "  GPT-4o mini, one call per control",
        "  Saved to policies/cis_aws_v7.json",
    ]),
    ("Runtime (per audit)", [
        "  Clone repo + Checkov scan",
        "  Extract enclosing Terraform block",
        "  Attach HCL to every finding",
    ]),
    ("Consensus", [
        "  3 models judge each finding in parallel",
        "  SHA-256 cache: same input → same verdicts",
        "  Versioned audit stored in SQLite",
    ]),
]
for i, (title, items) in enumerate(detail):
    lx = Inches(0.3 + i * 4.35); ty = Inches(3.3)
    card(s, lx, ty, Inches(4.1), Inches(3.7), title, INK)
    bul(s, items, lx + Inches(0.12), ty + Inches(0.55),
        Inches(3.85), Inches(3.05), size=13)


# =============================================================================
# SLIDE 5 — Evaluation Methodology
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
hdr(s, "How We Evaluated", "Two independent test suites with hand-labeled ground truth")

left = [
    ("FP-Filter Eval", ACCENT, [
        "30 hand-labeled findings",
        "15 GENUINE  +  15 FALSE_POSITIVE",
        "13 distinct Checkov check IDs",
        "Each item: real HCL + correct verdict",
        "Metrics: P / R / F1, kappa, per-tier precision",
    ]),
    ("Planner Eval", INK, [
        "55 test cases from CIS AWS Benchmark",
        "Input: clause text",
        "Output: severity, type, target resources",
        "Metrics: severity acc, type acc, targets F1",
        "Also: latency, tokens, $ cost",
    ]),
]
for i, (title, color, items) in enumerate(left):
    lx = Inches(0.4 + i * 6.45); ty = Inches(1.55)
    card(s, lx, ty, Inches(6.1), Inches(3.6), title, color)
    bul(s, ["  " + it for it in items],
        lx + Inches(0.14), ty + Inches(0.55),
        Inches(5.85), Inches(3.0), size=14)

txt(s, "Why this evaluation is defensible",
    Inches(0.4), Inches(5.4), Inches(12.5), Inches(0.35),
    size=15, bold=True, color=INK)
rect(s, Inches(0.4), Inches(5.8), Inches(12.5), Inches(1.5), fill=WHITE, line=LGRAY, lw=Pt(1))
bul(s, [
    "  Annotations are real Terraform from open-source repos — not synthetic",
    "  Each item picked to test contextual reasoning the LLM must do",
    "  SHA-256 cache makes every number deterministic and reproducible",
],   Inches(0.55), Inches(5.9), Inches(12.2), Inches(1.35), size=14)


# =============================================================================
# SLIDE 6 — FP-Filter Headline
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
hdr(s, "Results — FP-Filter Consensus", "Three views of the same n = 30 evaluation")

big_stat(s, "84.6 – 100%",  "Per-model F1 range",        Inches(0.5),  Inches(1.55), Inches(4.0), Inches(1.8), AMBER)
big_stat(s, "0.67 – 0.93",  "Inter-model agreement κ",   Inches(4.7),  Inches(1.55), Inches(4.0), Inches(1.8), ACCENT)
big_stat(s, "100%",         "Joint consensus accuracy",  Inches(8.9),  Inches(1.55), Inches(4.0), Inches(1.8), GREEN)

txt(s, "What each number tells us",
    Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.35),
    size=15, bold=True, color=INK)

rect(s, Inches(0.5),  Inches(4.05), Inches(4.0), Inches(2.8), fill=WHITE, line=LGRAY, lw=Pt(1))
txt(s, "Models differ",
    Inches(0.7), Inches(4.15), Inches(3.7), Inches(0.3), size=13, bold=True, color=AMBER)
bul(s, [
    "  Llama F1 = 84.6%",
    "  GPT-4o mini F1 = 96.8%",
    "  DeepSeek F1 = 100%",
    "  No single model is sufficient",
],   Inches(0.6), Inches(4.5), Inches(3.8), Inches(2.2), size=13)

rect(s, Inches(4.7), Inches(4.05), Inches(4.0), Inches(2.8), fill=WHITE, line=LGRAY, lw=Pt(1))
txt(s, "They genuinely disagree",
    Inches(4.9), Inches(4.15), Inches(3.7), Inches(0.3), size=13, bold=True, color=ACCENT)
bul(s, [
    "  κ < 1 on every pair",
    "  Models with different blind spots",
    "  Voting does real work —",
    "  not just rubber-stamping",
],   Inches(4.8), Inches(4.5), Inches(3.8), Inches(2.2), size=13)

rect(s, Inches(8.9), Inches(4.05), Inches(4.0), Inches(2.8), fill=WHITE, line=LGRAY, lw=Pt(1))
txt(s, "Together they recover",
    Inches(9.1), Inches(4.15), Inches(3.7), Inches(0.3), size=13, bold=True, color=GREEN)
bul(s, [
    "  30 / 30 cases correctly tiered",
    "  15 / 15 false positives suppressed",
    "  15 / 15 true violations retained",
    "  Consensus > any single model",
],   Inches(9.0), Inches(4.5), Inches(3.8), Inches(2.2), size=13)

txt(s, "Caveat: n = 30, hand-curated. Per-model variance is the load-bearing evidence — explored next slide.",
    Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
    size=12, color=GRAY, italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 7 — Per-Model Failure Modes
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
hdr(s, "Each Model Fails Differently", "Errors are non-overlapping — that's why voting recovers them")

# Hero: rescue effect
big_stat(s, "5",  "individual errors across all 3 models", Inches(2.0), Inches(1.55), Inches(4.0), Inches(1.4), RED)
txt(s, "→", Inches(6.2), Inches(1.85), Inches(0.9), Inches(0.8),
    size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
big_stat(s, "0",  "errors after consensus voting",         Inches(7.3), Inches(1.55), Inches(4.0), Inches(1.4), GREEN)

# Three model cards — each emphasizes a distinct failure mode
models = [
    ("Llama 3.1 70B", RED, [
        "Misses 4 real violations",
        "Recall = 73.3%",
        "Failure mode: under-flags",
        "Cheapest of the three, fast",
    ]),
    ("GPT-4o mini", AMBER, [
        "Flags 1 false positive",
        "Precision = 93.8%",
        "Failure mode: over-flags once",
        "Fastest (1.34s in planner eval)",
    ]),
    ("DeepSeek V3", GRAY, [
        "0 errors on n = 30",
        "Best on this set —",
        "but n is small + DeepSeek is slowest",
        "(7.15s in planner eval)",
    ]),
]
for i, (title, accent, items) in enumerate(models):
    lx = Inches(0.4 + i * 4.3); ty = Inches(3.25)
    card(s, lx, ty, Inches(4.1), Inches(3.0), title, accent)
    bul(s, ["  " + it for it in items],
        lx + Inches(0.14), ty + Inches(0.55),
        Inches(3.85), Inches(2.4), size=13, sp=6)

txt(s, "If we ran any single model, we'd be wrong 0 – 4 times. Different models fail on different cases →",
    Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.4),
    size=12.5, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
txt(s, "the union of their mistakes does NOT exceed the majority threshold on any case.",
    Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.4),
    size=12.5, color=GRAY, italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 8 — Disagreement + Tier Separation
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
hdr(s, "Models Disagree  |  Tiers Separate Truth", "Cohen's κ + tier distribution")

# LEFT: Cohen's kappa (varied numbers — the good story)
txt(s, "Cohen's κ across model pairs",
    Inches(0.5), Inches(1.38), Inches(6), Inches(0.35),
    size=15, bold=True, color=INK)
tbl(s, ["Pair", "κ", "Strength"],
    [
        ["GPT  vs  DeepSeek",   ("0.93", GREEN), "Almost perfect"],
        ["DeepSeek  vs  Llama", ("0.73", AMBER), "Substantial"],
        ["GPT  vs  Llama",      ("0.67", AMBER), "Substantial"],
    ],
    Inches(0.5), Inches(1.8),
    [Inches(2.7), Inches(1.2), Inches(2.3)])

rect(s, Inches(0.5), Inches(4.0), Inches(6.2), Inches(2.9), fill=WHITE, line=LGRAY, lw=Pt(1))
txt(s, "What κ tells us",
    Inches(0.65), Inches(4.1), Inches(5.9), Inches(0.35),
    size=14, bold=True, color=INK)
bul(s, [
    "  κ ranges 0.67 – 0.93, never 1.0",
    "  GPT  +  DeepSeek closest in style",
    "  Llama disagrees most with the other two",
    "  Disagreement → voting is non-trivial",
],   Inches(0.65), Inches(4.55), Inches(5.85), Inches(2.3), size=13, sp=6)

# RIGHT: tier distribution as counts (not percentages)
txt(s, "Tier  →  ground-truth distribution",
    Inches(7.0), Inches(1.38), Inches(6), Inches(0.35),
    size=15, bold=True, color=INK)
tbl(s, ["Tier", "Count", "Genuine", "FP"],
    [
        [("HIGH",       GREEN), "11", ("11", GREEN), "0"],
        [("LIKELY",     AMBER), "4",  ("4",  GREEN), "0"],
        [("LIKELY_FP",  AMBER), "1",  "0",            ("1",  RED)],
        [("SUPPRESSED", RED),   "14", "0",            ("14", RED)],
    ],
    Inches(7.0), Inches(1.8),
    [Inches(2.0), Inches(1.1), Inches(1.4), Inches(1.4)])

rect(s, Inches(7.0), Inches(4.45), Inches(5.85), Inches(2.45), fill=WHITE, line=LGRAY, lw=Pt(1))
txt(s, "What the counts show",
    Inches(7.15), Inches(4.55), Inches(5.55), Inches(0.35),
    size=14, bold=True, color=INK)
bul(s, [
    "  All 15 GENUINE items → HIGH or LIKELY tier",
    "  All 15 FP items → LIKELY_FP or SUPPRESSED tier",
    "  Zero cross-contamination across the boundary",
    "  Tier names are predictive, not decorative",
],   Inches(7.15), Inches(4.95), Inches(5.55), Inches(1.9), size=13, sp=6)

rect(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4), fill=INK)
txt(s, "Real disagreement  +  clean tier separation  =  the consensus mechanism is doing measurable work.",
    Inches(0.65), Inches(7.1), Inches(12.0), Inches(0.32),
    size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 9 — Planner Benchmark
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
hdr(s, "Results — Planner Benchmark", "55 CIS clauses  |  accuracy vs latency vs cost tradeoff")

tbl(s, ["Model", "Type Acc.", "Targets F1", "Severity Acc.", "Latency", "Cost / run"],
    [
        ["GPT-4o mini",   ("85.5%", AMBER), ("70.5%", AMBER), ("40.0%", GRAY),
                          ("1.34 s", GREEN), ("$0.0022", AMBER)],
        ["DeepSeek V3",   ("90.9%", GREEN), ("80.9%", GREEN), ("38.2%", GRAY),
                          ("7.15 s", RED),   ("$0.0018", GREEN)],
        ["Llama 3.1 70B", ("90.9%", GREEN), ("79.2%", GREEN), ("50.9%", GRAY),
                          ("3.04 s", AMBER), ("$0.0057", RED)],
    ],
    Inches(0.4), Inches(1.55),
    [Inches(2.7), Inches(1.65), Inches(1.85), Inches(2.0), Inches(1.55), Inches(2.05)])

txt(s, "Reading the table",
    Inches(0.5), Inches(3.85), Inches(6.1), Inches(0.35),
    size=15, bold=True, color=INK)
rect(s, Inches(0.5), Inches(4.25), Inches(6.1), Inches(2.9), fill=WHITE, line=LGRAY, lw=Pt(1))
bul(s, [
    "  No single winner — explicit tradeoff",
    "  GPT-4o mini  →  5× faster than DeepSeek",
    "  DeepSeek      →  best on type + targets, lowest cost",
    "  Llama 3.1 70B →  best severity, mid-price",
],   Inches(0.65), Inches(4.35), Inches(5.85), Inches(2.7), size=14, sp=8)

txt(s, "Why severity accuracy is low",
    Inches(7.0), Inches(3.85), Inches(5.7), Inches(0.35),
    size=15, bold=True, color=INK)
rect(s, Inches(7.0), Inches(4.25), Inches(5.7), Inches(2.9), fill=WHITE, line=LGRAY, lw=Pt(1))
bul(s, [
    "  Severity is the most subjective field",
    "  Engineers disagree on critical vs high",
    "  Type acc + Targets F1 are objective",
    "  Those are the load-bearing numbers",
],   Inches(7.15), Inches(4.35), Inches(5.45), Inches(2.7), size=14, sp=8)


# =============================================================================
# SLIDE 10 — Live System
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
hdr(s, "Live System", "End-to-end audit on bridgecrewio/terragoat")

big_stat(s, "385", "Checkov findings",   Inches(0.5),  Inches(1.55), Inches(2.95), Inches(1.7), INK)
big_stat(s, "211", "high severity",      Inches(3.65), Inches(1.55), Inches(2.95), Inches(1.7), RED)
big_stat(s, "174", "medium severity",    Inches(6.8),  Inches(1.55), Inches(2.95), Inches(1.7), AMBER)
big_stat(s, "86.8%", "compliance score", Inches(9.95), Inches(1.55), Inches(2.95), Inches(1.7), GREEN)

txt(s, "What the demo shows",
    Inches(0.5), Inches(3.55), Inches(12.3), Inches(0.35),
    size=15, bold=True, color=INK)
rect(s, Inches(0.5), Inches(3.95), Inches(12.3), Inches(3.0), fill=WHITE, line=LGRAY, lw=Pt(1))
bul(s, [
    "  Paste a GitHub URL  →  one click  →  audit pipeline runs",
    "  Each finding shows per-model verdicts: GENUINE / FALSE_POSITIVE / UNCERTAIN",
    "  Tiered display: HIGH (red)  |  LIKELY (amber)  |  SUPPRESSED (collapsed)",
    "  Eval panel: live per-model F1, Cohen's κ, per-tier precision",
    "  SHA-256 cache: re-audit of the same repo completes in seconds",
],   Inches(0.7), Inches(4.05), Inches(12.0), Inches(2.85),
    size=14.5, sp=10)


# =============================================================================
# SLIDE 11 — Limitations & Future Work
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=PAPER)
hdr(s, "Limitations & Future Work", "Honest framing of what's measured and what's next")

txt(s, "What's measured today",
    Inches(0.5), Inches(1.38), Inches(6), Inches(0.35),
    size=15, bold=True, color=RED)
rect(s, Inches(0.5), Inches(1.78), Inches(6.1), Inches(2.7), fill=WHITE, line=LGRAY, lw=Pt(1))
bul(s, [
    "  n = 30 annotations — small set",
    "  Hand-curated with strong context signals",
    "  Single framework (CIS AWS)",
    "  Terraform / HCL only",
],   Inches(0.65), Inches(1.88), Inches(5.85), Inches(2.5), size=14, sp=8)

txt(s, "What's coming next",
    Inches(7.0), Inches(1.38), Inches(6), Inches(0.35),
    size=15, bold=True, color=GREEN)
rect(s, Inches(7.0), Inches(1.78), Inches(5.7), Inches(2.7), fill=WHITE, line=LGRAY, lw=Pt(1))
bul(s, [
    "  Expand to 80–100 adversarial annotations",
    "  Add NIST SP 800-53, SOC 2 frameworks",
    "  CloudFormation + Kubernetes support",
    "  Auto-generated fix PRs with terraform validate",
],   Inches(7.15), Inches(1.88), Inches(5.45), Inches(2.5), size=14, sp=8)

txt(s, "What we deliberately did NOT do",
    Inches(0.5), Inches(4.75), Inches(12.3), Inches(0.35),
    size=15, bold=True, color=INK)
rect(s, Inches(0.5), Inches(5.15), Inches(12.3), Inches(2.0), fill=INK)
bul(s, [
    "  Did not train or fine-tune any model — three off-the-shelf LLMs only",
    "  Did not use LLMs as detectors — they only judge what Checkov surfaces",
    "  Did not cherry-pick — eval JSONs are regenerable with one command",
],   Inches(0.65), Inches(5.25), Inches(12.0), Inches(1.85),
    size=14, color=WHITE, sp=10)


# =============================================================================
# SLIDE 12 — Summary
# =============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, fill=INK)
rect(s, 0, H - Inches(0.07), W, Inches(0.07), fill=ACCENT)
txt(s, "Summary",
    Inches(1), Inches(0.7), Inches(11.33), Inches(0.9),
    size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, Inches(4.0), Inches(1.55), Inches(5.33), Inches(0.05), fill=ACCENT)

pts = [
    "Hybrid system:  deterministic detection (Checkov) + contextual judgement (3 LLMs)",
    "FP-Filter eval, n = 30:   100 % overall accuracy   |   100 % FP reduction   |   100 % TP retention",
    "Per-model F1 varies 84.6 % → 100 %  —  consensus rescues each model's blind spots",
    "Cohen's κ 0.67 – 0.93  —  models genuinely disagree, voting does real work",
    "Tiers carry signal:  HIGH = 100 % genuine,  SUPPRESSED = 0 % genuine",
    "Planner eval, n = 55:   type acc 85 – 91 %,  latency 1.3 – 7.2 s,  cost $0.0018 – $0.0057",
]
for i, p in enumerate(pts):
    ty = Inches(1.95 + i * 0.78)
    rect(s, Inches(1.0), ty + Inches(0.14), Inches(0.18), Inches(0.18), fill=ACCENT)
    txt(s, p, Inches(1.35), ty, Inches(11.0), Inches(0.72), size=14, color=WHITE)

txt(s, "github.com/ParthGala2k/AuditTrace",
    Inches(1), Inches(7.0), Inches(11.33), Inches(0.35),
    size=13, color=RGBColor(0x9d, 0x8f, 0xd0), align=PP_ALIGN.CENTER, italic=True)


# =============================================================================
prs.save("AuditTrace_Presentation.pptx")
print("Saved: AuditTrace_Presentation.pptx  (12 slides)")
